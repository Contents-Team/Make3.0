"""Extract text from electrical wiring storyboard PPTX."""
import sys
from pathlib import Path
from pptx import Presentation

PATH = Path(r"C:\Users\VIRNECT\Downloads\work\01_Other\Make3.0") / "[작성중] [폴리텍신기술교육원]_전기내선공사 가상 실습_SB_v1.0.pptx"

prs = Presentation(PATH)
import io, os
OUT = io.StringIO()
import builtins as _b
def p(*a, **k):
    _b.print(*a, **k, file=OUT)

p(f"Slides: {len(prs.slides)}")
p(f"Size: {prs.slide_width} x {prs.slide_height}")
p("=" * 80)

for i, slide in enumerate(prs.slides, start=1):
    p(f"\n=== SLIDE {i} ===")
    # title if present
    if slide.shapes.title and slide.shapes.title.text:
        p(f"TITLE: {slide.shapes.title.text}")
    # all text frames
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                txt = para.text.strip()
                if txt:
                    p(f"  {txt}")
        # tables
        if shape.has_table:
            p("  [TABLE]")
            for row in shape.table.rows:
                cells = [cell.text.strip().replace("\n", " | ") for cell in row.cells]
                p("    | " + " | ".join(cells) + " |")
    # notes
    if slide.has_notes_slide:
        notes = slide.notes_slide.notes_text_frame.text.strip()
        if notes:
            p(f"  [NOTES] {notes[:500]}")

out_path = os.path.join(os.environ.get("TEMP", "/tmp"), "pptx_dump.txt")
with open(out_path, "w", encoding="utf-8") as f:
    f.write(OUT.getvalue())
print(f"Wrote {len(OUT.getvalue())} chars to {out_path}")
